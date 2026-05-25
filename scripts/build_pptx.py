#!/usr/bin/env python3
"""Build the NOC Whisperer capstone PowerPoint deck (May 30, 2026) — 13 slides."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

PROJECT_ROOT = Path(__file__).resolve().parents[1]
OUTPUT_PATH = PROJECT_ROOT / "docs" / "NOC_Whisperer_Capstone.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BG_PRIMARY = "1E1E2E"
BG_CARD = "2D2D42"
BG_CARD_ALT = "252535"
TEXT_PRIMARY = "FFFFFF"
TEXT_MUTED = "A0A0B8"
LLM_BLUE = "4472C4"
RULE_GREEN = "70AD47"
MCP_ORANGE = "ED7D31"
DATA_GRAY = "A5A5A5"
GOLD = "FFD700"
ALERT_RED = "FF4444"
CONFIRMED_TEAL = "00C897"
OTEL_PURPLE = "9B59B6"
ARROW_GRAY = "555570"
BORDER_MUTED = "3D3D55"


def rgb(hex_value: str) -> RGBColor:
    """Convert 6-digit hex to RGBColor."""
    h = hex_value.strip().lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def new_presentation() -> Presentation:
    """Create a blank widescreen presentation."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs: Presentation):
    """Add a blank slide with dark background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(BG_PRIMARY)
    return slide


def set_shape_fill(shape, hex_fill: str) -> None:
    """Apply solid fill."""
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(hex_fill)


def set_shape_border(shape, hex_color: str, width_pt: float = 1.0) -> None:
    """Apply line border."""
    shape.line.color.rgb = rgb(hex_color)
    shape.line.width = Pt(width_pt)


def no_border(shape) -> None:
    """Remove shape outline."""
    shape.line.fill.background()


def add_rounded_rect(
    slide,
    left,
    top,
    width,
    height,
    *,
    fill_hex: str = BG_CARD,
    border_hex: str | None = None,
    border_pt: float = 1.0,
):
    """Add a rounded rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    set_shape_fill(shape, fill_hex)
    if border_hex:
        set_shape_border(shape, border_hex, border_pt)
    else:
        no_border(shape)
    return shape


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    font_size: int = 15,
    bold: bool = False,
    italic: bool = False,
    color_hex: str = TEXT_PRIMARY,
    font_name: str = "Calibri",
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    word_wrap: bool = True,
):
    """Add a single-paragraph text box."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = word_wrap
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = rgb(color_hex)
    return box


def add_slide_title(slide, text: str) -> None:
    """Left-aligned section title (22pt bold)."""
    add_textbox(
        slide,
        Inches(0.4),
        Inches(0.35),
        Inches(12.5),
        Inches(0.55),
        text,
        font_size=22,
        bold=True,
    )


def add_subtitle(slide, text: str, *, top=Inches(0.95), size: int = 13) -> None:
    """Muted subtitle under title."""
    add_textbox(
        slide,
        Inches(0.4),
        top,
        Inches(12.5),
        Inches(0.4),
        text,
        font_size=size,
        color_hex=TEXT_MUTED,
    )


def add_horizontal_arrow(slide, left, top, width=Inches(0.25)) -> None:
    """Thin horizontal arrow with triangle tip."""
    bar_h = Inches(0.04)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, bar_h)
    set_shape_fill(bar, ARROW_GRAY)
    no_border(bar)
    tip = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        left + width,
        top - Inches(0.03),
        Inches(0.12),
        Inches(0.1),
    )
    set_shape_fill(tip, ARROW_GRAY)
    no_border(tip)
    tip.rotation = 90.0


def add_labeled_arrow_right(
    slide,
    left,
    top,
    length: float,
    label: str,
    *,
    arrow_color: str = DATA_GRAY,
    label_color: str = TEXT_MUTED,
    italic: bool = True,
    dashed: bool = False,
) -> None:
    """Horizontal arrow with label above."""
    w = Inches(length)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, Inches(0.02))
    set_shape_fill(bar, arrow_color)
    if dashed:
        bar.line.color.rgb = rgb(arrow_color)
        bar.line.width = Pt(1)
        bar.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    else:
        no_border(bar)
    tip = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        left + w,
        top - Inches(0.02),
        Inches(0.1),
        Inches(0.08),
    )
    set_shape_fill(tip, arrow_color)
    no_border(tip)
    tip.rotation = 90.0
    add_textbox(
        slide,
        left,
        top - Inches(0.22),
        w + Inches(0.15),
        Inches(0.2),
        label,
        font_size=10,
        italic=italic,
        color_hex=label_color,
        align=PP_ALIGN.CENTER,
    )


def add_vertical_arrow(slide, left, top, height=Inches(0.4), color_hex: str = MCP_ORANGE) -> None:
    """Downward arrow."""
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, Inches(0.28), height)
    set_shape_fill(arrow, color_hex)
    no_border(arrow)


def add_filled_circle(slide, left, top, diameter, color_hex: str) -> None:
    """Small filled circle marker."""
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
    set_shape_fill(c, color_hex)
    no_border(c)


def add_left_accent_row(
    slide,
    left,
    top,
    width,
    height,
    accent_hex: str,
    body: str,
    *,
    body_size: int = 13,
) -> None:
    """Text row with colored left bar."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, Inches(0.06), height
    )
    set_shape_fill(bar, accent_hex)
    no_border(bar)
    add_textbox(
        slide,
        left + Inches(0.12),
        top,
        width - Inches(0.12),
        height,
        body,
        font_size=body_size,
    )


# ─── Slide builders ───────────────────────────────────────────────────────────


def build_slide_1_title(prs: Presentation) -> None:
    """Slide 1 — Title."""
    slide = blank_slide(prs)
    band_w = Inches(9.0)
    band = add_rounded_rect(
        slide,
        (SLIDE_W - band_w) / 2,
        Inches(3.2),
        band_w,
        Inches(0.08),
        fill_hex=LLM_BLUE,
    )
    no_border(band)
    add_textbox(
        slide,
        Inches(0.5),
        Inches(2.1),
        Inches(12.3),
        Inches(0.9),
        "NOC WHISPERER",
        font_size=52,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Inches(0.8),
        Inches(3.5),
        Inches(11.7),
        Inches(0.9),
        "Real-time multi-agent alert correlation\n"
        "for distributed infrastructure",
        font_size=20,
        color_hex=TEXT_MUTED,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Inches(0.5),
        Inches(6.6),
        Inches(12.3),
        Inches(0.4),
        "H2M Agents  |  May 30, 2026",
        font_size=14,
        color_hex=TEXT_MUTED,
        align=PP_ALIGN.CENTER,
    )


def build_slide_2_problem(prs: Presentation) -> None:
    """Slide 2 — The Problem."""
    slide = blank_slide(prs)
    add_slide_title(slide, "The NOC Alert Flood Problem")

    panel = add_rounded_rect(
        slide,
        Inches(0.45),
        Inches(1.3),
        Inches(12.4),
        Inches(0.9),
        fill_hex=BG_CARD,
        border_hex=OTEL_PURPLE,
        border_pt=1.5,
    )
    add_filled_circle(slide, Inches(0.6), Inches(1.55), Inches(0.2), OTEL_PURPLE)
    add_textbox(
        slide,
        Inches(0.95),
        Inches(1.42),
        Inches(11.7),
        Inches(0.75),
        "OpenTelemetry is an open-source observability framework that instruments "
        "applications and emits metrics, traces, and logs — the signals this system "
        "monitors in real time.",
        font_size=13,
    )

    bullets = [
        "A single infrastructure failure cascades across 7+ services "
        "— each generating its own alert",
        "L1 operators manually correlate dozens of noisy alerts",
        "Root cause identification takes 15–30 minutes",
        "Alert fatigue leads to missed signals",
    ]
    y0 = Inches(2.4)
    for i, bullet in enumerate(bullets):
        y = y0 + i * Inches(0.82)
        add_filled_circle(slide, Inches(0.6), y + Inches(0.06), Inches(0.25), ALERT_RED)
        add_textbox(slide, Inches(1.0), y, Inches(11.8), Inches(0.75), bullet, font_size=15)

    callout = add_rounded_rect(
        slide,
        Inches(1.15),
        Inches(5.8),
        Inches(11.0),
        Inches(1.1),
        fill_hex=BG_CARD,
        border_hex=GOLD,
        border_pt=2.0,
    )
    tf = callout.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = (
        "Can a multi-agent AI system correlate alerts in real time, "
        "identify root cause, and generate actionable advisories — automatically?"
    )
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Calibri"
    p.font.size = Pt(15)
    p.font.italic = True
    p.font.color.rgb = rgb(GOLD)


def _two_line_box(
    slide,
    left,
    top,
    width,
    height,
    line1: str,
    line2: str,
    *,
    border_hex: str,
    line1_bold: bool = True,
) -> None:
    """Rounded box with title + subtitle lines."""
    box = add_rounded_rect(
        slide, left, top, width, height, fill_hex=BG_CARD, border_hex=border_hex
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_top = Inches(0.06)
    p1 = tf.paragraphs[0]
    p1.text = line1
    p1.font.name = "Calibri"
    p1.font.size = Pt(13)
    p1.font.bold = line1_bold
    p1.font.color.rgb = rgb(TEXT_PRIMARY)
    p2 = tf.add_paragraph()
    p2.text = line2
    p2.font.size = Pt(11)
    p2.font.color.rgb = rgb(TEXT_MUTED)


def _agent_pill(
    slide,
    left,
    top,
    width,
    height,
    name: str,
    type_label: str,
    border_hex: str,
) -> None:
    """Agent pill with name left and type right."""
    pill = add_rounded_rect(
        slide, left, top, width, height, fill_hex=BG_CARD, border_hex=border_hex
    )
    tf = pill.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = f"{name}                    {type_label}"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = rgb(TEXT_PRIMARY)


def build_slide_3_system_context(prs: Presentation) -> None:
    """Slide 3 — System Context diagram."""
    slide = blank_slide(prs)
    add_slide_title(slide, "System Context")
    add_subtitle(
        slide,
        "What flows in, what flows out, and who is outside the system boundary",
    )

    add_textbox(
        slide,
        Inches(0.4),
        Inches(1.4),
        Inches(3.2),
        Inches(0.25),
        "MONITORED INFRASTRUCTURE",
        font_size=11,
        bold=True,
        color_hex=TEXT_MUTED,
        align=PP_ALIGN.CENTER,
    )

    sources = [
        ("OTel Demo Stack", "applications + services", OTEL_PURPLE),
        ("Prometheus", "metrics + thresholds", DATA_GRAY),
        ("Jaeger", "distributed traces", DATA_GRAY),
        ("Node Exporter", "host-level metrics", DATA_GRAY),
    ]
    sy = Inches(1.8)
    box_h = Inches(0.7)
    gap = Inches(0.15)
    for i, (title, sub, border) in enumerate(sources):
        y = sy + i * (box_h + gap)
        _two_line_box(
            slide, Inches(0.5), y, Inches(3.0), box_h, title, sub, border_hex=border
        )

    # Zone 2 — system boundary
    zx, zy, zw, zh = Inches(3.9), Inches(1.3), Inches(5.5), Inches(5.2)
    boundary = add_rounded_rect(
        slide, zx, zy, zw, zh, fill_hex=BG_CARD_ALT, border_hex=LLM_BLUE, border_pt=2.5
    )
    add_textbox(
        slide,
        zx + Inches(0.3),
        Inches(1.55),
        zw - Inches(0.6),
        Inches(0.4),
        "NOC WHISPERER",
        font_size=18,
        bold=True,
        color_hex=LLM_BLUE,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        zx + Inches(0.3),
        Inches(1.95),
        zw - Inches(0.6),
        Inches(0.35),
        "Multi-Agent Alert Correlation System",
        font_size=12,
        color_hex=TEXT_MUTED,
        align=PP_ALIGN.CENTER,
    )

    agents = [
        ("NormalizerAgent", "LLM", LLM_BLUE),
        ("TriageAgent", "rule-based", RULE_GREEN),
        ("CorrelationAgent", "DSPy", LLM_BLUE),
        ("CommunicationsAgent", "LLM", LLM_BLUE),
    ]
    px = zx + Inches(0.5)
    pw = Inches(4.5)
    ph = Inches(0.55)
    pgap = Inches(0.12)
    py0 = Inches(2.4)
    for i, (name, typ, border) in enumerate(agents):
        _agent_pill(slide, px, py0 + i * (ph + pgap), pw, ph, name, typ, border)

    orch = add_rounded_rect(
        slide,
        px,
        Inches(5.05),
        pw,
        Inches(0.4),
        fill_hex=BG_CARD,
        border_hex=LLM_BLUE,
    )
    tf = orch.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "ADK LlmAgent Orchestrator — 15s cycle"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(11)
    p.font.color.rgb = rgb(TEXT_PRIMARY)

    # Input arrows zone 1 → 2
    arrow_labels = [
        "metrics & threshold alerts",
        "distributed traces",
        "host metrics",
    ]
    for i, label in enumerate(arrow_labels):
        y = sy + i * (box_h + gap) + box_h / 2 - Inches(0.02)
        add_labeled_arrow_right(slide, Inches(3.55), y, 0.28, label)

    # Zone 3
    add_textbox(
        slide,
        Inches(9.7),
        Inches(1.4),
        Inches(3.2),
        Inches(0.25),
        "NOC OPERATORS",
        font_size=11,
        bold=True,
        color_hex=TEXT_MUTED,
        align=PP_ALIGN.CENTER,
    )
    _two_line_box(
        slide,
        Inches(9.85),
        Inches(1.8),
        Inches(3.0),
        Inches(1.0),
        "NOC Team",
        "on-call engineers",
        border_hex=CONFIRMED_TEAL,
    )

    out_labels = ["NOC Advisory", "Incident Board", "SERVICE RESTORED"]
    for i, label in enumerate(out_labels):
        y = Inches(2.2) + i * Inches(0.55)
        add_labeled_arrow_right(
            slide,
            Inches(9.45),
            y,
            0.22,
            label,
            arrow_color=CONFIRMED_TEAL,
            label_color=CONFIRMED_TEAL,
            italic=False,
        )

    add_textbox(
        slide,
        Inches(0.5),
        Inches(6.9),
        Inches(12.3),
        Inches(0.45),
        "System ingests observability signals autonomously — "
        "operators receive correlated advisories, not raw alerts",
        font_size=11,
        italic=True,
        color_hex=TEXT_MUTED,
        align=PP_ALIGN.CENTER,
    )


def build_slide_4_architecture(prs: Presentation) -> None:
    """Slide 4 — System Architecture (internal)."""
    slide = blank_slide(prs)
    add_slide_title(slide, "System Architecture")

    columns = [
        ("DATA SOURCES", DATA_GRAY, ["Prometheus", "Jaeger", "Node Exporter", "OTel Demo Stack"]),
        ("MCP TOOL LAYER", MCP_ORANGE, ["PrometheusMCP", "JaegerMCP", "TopologyMCP", "NodeExporterMCP"]),
        (
            "AGENT PIPELINE",
            LLM_BLUE,
            [
                ("NormalizerAgent", LLM_BLUE, "LLM fine-tuned"),
                ("TriageAgent", RULE_GREEN, "rule-based"),
                ("CorrelationAgent", LLM_BLUE, "LLM DSPy"),
                ("CommunicationsAgent", LLM_BLUE, "LLM fine-tuned"),
            ],
        ),
        (
            "NOC DASHBOARD",
            CONFIRMED_TEAL,
            ["RAW ALERT STREAM", "INCIDENT BOARD", "NOC ADVISORY", "SQLite IncidentStore"],
        ),
    ]
    col_x = [Inches(0.4), Inches(3.4), Inches(6.4), Inches(9.4)]
    col_w = Inches(2.6)
    header_y = Inches(1.35)
    header_h = Inches(0.4)
    box_y0 = Inches(1.9)
    box_h = Inches(0.75)
    box_gap = Inches(0.12)
    arrow_w = Inches(0.25)

    for col_idx, (header, accent, items) in enumerate(columns):
        x = col_x[col_idx]
        hdr = add_rounded_rect(slide, x, header_y, col_w, header_h, fill_hex=accent)
        no_border(hdr)
        tf = hdr.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = header
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = rgb(TEXT_PRIMARY)

        for row_idx, item in enumerate(items):
            y = box_y0 + row_idx * (box_h + box_gap)
            if col_idx == 2:
                name, border_c, caption = item  # type: ignore[misc]
                box = add_rounded_rect(
                    slide, x, y, col_w, box_h, fill_hex=BG_CARD, border_hex=border_c, border_pt=1.5
                )
                tf_b = box.text_frame
                tf_b.word_wrap = True
                tf_b.margin_left = Inches(0.08)
                tf_b.margin_top = Inches(0.05)
                p1 = tf_b.paragraphs[0]
                p1.text = name
                p1.font.size = Pt(12)
                p1.font.bold = True
                p1.font.color.rgb = rgb(TEXT_PRIMARY)
                p2 = tf_b.add_paragraph()
                p2.text = caption
                p2.font.size = Pt(10)
                p2.font.color.rgb = rgb(TEXT_MUTED)
            else:
                label = str(item)
                box = add_rounded_rect(
                    slide, x, y, col_w, box_h, fill_hex=BG_CARD, border_hex=accent
                )
                tf_b = box.text_frame
                tf_b.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = tf_b.paragraphs[0]
                p.text = label
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(12)
                p.font.color.rgb = rgb(TEXT_PRIMARY)

            if col_idx < 3:
                ay = y + box_h / 2 - Inches(0.02)
                add_horizontal_arrow(slide, x + col_w + Inches(0.05), ay, width=arrow_w)

    add_textbox(
        slide,
        Inches(0.5),
        Inches(6.75),
        Inches(12.3),
        Inches(0.4),
        "Jaeger / NodeExporter MCPs in codebase; not wired into ADK demo loop",
        font_size=11,
        italic=True,
        color_hex=TEXT_MUTED,
        align=PP_ALIGN.CENTER,
    )


def build_slide_5_tool_stack(prs: Presentation) -> None:
    """Slide 5 — Two Layers of Tools."""
    slide = blank_slide(prs)
    add_slide_title(slide, "Two Layers of Tools")

    panel_h = Inches(4.8)
    panel_y = Inches(1.3)
    panel_w = Inches(5.8)

    hdr1 = add_rounded_rect(
        slide, Inches(0.5), panel_y, panel_w, Inches(0.5), fill_hex=LLM_BLUE
    )
    no_border(hdr1)
    tf = hdr1.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "LAYER 1 — ADK FunctionTools"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = rgb(TEXT_PRIMARY)

    add_rounded_rect(slide, Inches(0.5), panel_y, panel_w, panel_h, fill_hex=BG_CARD)
    add_textbox(
        slide, Inches(0.65), panel_y + Inches(0.52), panel_w, Inches(0.3),
        "(agents/adk_tools/)", font_size=12, color_hex=TEXT_MUTED,
    )
    add_textbox(
        slide, Inches(0.65), panel_y + Inches(0.85), panel_w, Inches(0.35),
        "What the LLM agent sees and calls:", font_size=13, italic=True, color_hex=TEXT_MUTED,
    )
    tools = [
        ("get_active_alerts()", "normalize_alert()"),
        ("route_alert()", "correlate_alert()"),
        ("generate_advisory()", "check_open_incidents()"),
        ("check_service_health()", "close_incident()"),
    ]
    gy = panel_y + Inches(1.25)
    for i, (a, b) in enumerate(tools):
        y = gy + i * Inches(0.55)
        add_textbox(slide, Inches(0.65), y, Inches(2.6), Inches(0.5), a, font_size=13,
                    color_hex=LLM_BLUE, font_name="Consolas")
        add_textbox(slide, Inches(3.35), y, Inches(2.6), Inches(0.5), b, font_size=13,
                    color_hex=LLM_BLUE, font_name="Consolas")

    add_vertical_arrow(slide, Inches(6.35), Inches(3.35))
    add_textbox(slide, Inches(6.75), Inches(3.45), Inches(0.8), Inches(0.4),
                "calls", font_size=12, color_hex=MCP_ORANGE, bold=True)

    rx = Inches(6.8)
    hdr2 = add_rounded_rect(slide, rx, panel_y, panel_w, Inches(0.5), fill_hex=MCP_ORANGE)
    no_border(hdr2)
    tf2 = hdr2.text_frame
    tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
    p2 = tf2.paragraphs[0]
    p2.text = "LAYER 2 — MCP Tool Classes"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(14)
    p2.font.bold = True
    p2.font.color.rgb = rgb(TEXT_PRIMARY)

    add_rounded_rect(slide, rx, panel_y, panel_w, panel_h, fill_hex=BG_CARD)
    add_textbox(slide, rx + Inches(0.15), panel_y + Inches(0.52), panel_w, Inches(0.3),
                "(mcp_tools/)", font_size=12, color_hex=TEXT_MUTED)
    add_textbox(slide, rx + Inches(0.15), panel_y + Inches(0.85), panel_w, Inches(0.35),
                "Handle infrastructure directly:", font_size=13, italic=True, color_hex=TEXT_MUTED)

    mcps = [
        ("PrometheusMCP", "Prometheus HTTP API (10.0.50.60:9090)"),
        ("JaegerMCP", "Jaeger trace API"),
        ("TopologyMCP", "otel_demo_graph.json"),
        ("NodeExporterMCP", "Node Exporter metrics"),
    ]
    for i, (cls, target) in enumerate(mcps):
        y = panel_y + Inches(1.35) + i * Inches(0.72)
        add_textbox(slide, rx + Inches(0.15), y, Inches(2.0), Inches(0.4), cls,
                    font_size=13, color_hex=MCP_ORANGE, font_name="Consolas")
        add_textbox(slide, rx + Inches(2.1), y, Inches(0.25), Inches(0.4), "→",
                    font_size=13, color_hex=TEXT_MUTED, font_name="Consolas")
        add_textbox(slide, rx + Inches(2.35), y, Inches(3.2), Inches(0.55), target,
                    font_size=12, color_hex=MCP_ORANGE, font_name="Consolas")

    add_textbox(
        slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.45),
        "The agent reasons about tools. The tools handle infrastructure.",
        font_size=13, italic=True, color_hex=TEXT_MUTED, align=PP_ALIGN.CENTER,
    )


def _step_card(slide, left, top, num: str, title: str, body: str) -> None:
    """Numbered step card."""
    add_rounded_rect(slide, left, top, Inches(5.8), Inches(1.25), fill_hex=BG_CARD)
    add_textbox(slide, left + Inches(0.12), top + Inches(0.06), Inches(0.45), Inches(0.4),
                num, font_size=20, bold=True, color_hex=LLM_BLUE)
    add_textbox(slide, left + Inches(0.5), top + Inches(0.06), Inches(5.1), Inches(0.38),
                title, font_size=14, bold=True, font_name="Consolas")
    add_textbox(slide, left + Inches(0.5), top + Inches(0.46), Inches(5.1), Inches(0.72),
                body, font_size=13, color_hex=TEXT_MUTED)


def build_slide_6_how_it_works(prs: Presentation) -> None:
    """Slide 6 — 6 Steps."""
    slide = blank_slide(prs)
    add_slide_title(slide, "From Raw Alert to NOC Advisory — 6 Steps")

    c1, c2 = Inches(0.5), Inches(6.9)
    ys = [Inches(1.35), Inches(2.75), Inches(4.15)]
    steps = [
        ("①", "get_active_alerts()", "Pull threshold breaches from Prometheus MCP"),
        ("②", "normalize_alert()", "Raw metric → canonical CanonicalAlert"),
        ("③", "route_alert()", "Rule-based topology + temporal proximity →\nnew incident or append to existing"),
        ("④", "correlate_alert()", "DSPy LLM reasons over alert cluster +\ntopology → root cause + confidence"),
        ("⑤", "generate_advisory()", "Fine-tuned model writes NOC advisory\n(preliminary → confirmed → resolution)"),
        ("⑥", "check_service_health()", "Prometheus confirms recovery →\nclose_incident() → SERVICE RESTORED"),
    ]
    positions = [(c1, ys[0]), (c2, ys[0]), (c1, ys[1]), (c2, ys[1]), (c1, ys[2]), (c2, ys[2])]
    for (num, title, body), (left, top) in zip(steps, positions):
        _step_card(slide, left, top, num, title, body)

    strip = add_rounded_rect(slide, Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.45), fill_hex=BG_CARD)
    no_border(strip)
    tf = strip.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "Cycle time: 15s  |  CONFIRMED when confidence > 0.85 AND alert_count ≥ 2"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(13)
    p.font.color.rgb = rgb(TEXT_PRIMARY)


def build_slide_7_models(prs: Presentation) -> None:
    """Slide 7 — Fine-tuned Models + RLVR discovery."""
    slide = blank_slide(prs)
    add_slide_title(slide, "Three LLM-Powered Agents")
    add_textbox(
        slide, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.35),
        "BASE MODEL: Qwen2.5-7B-Instruct  |  LoRA adapters ~10MB each",
        font_size=13, color_hex=TEXT_MUTED,
    )

    rows, cols = 4, 5
    tbl = slide.shapes.add_table(rows, cols, Inches(0.45), Inches(1.45), Inches(12.4), Inches(1.4)).table
    headers = ["Agent", "Training", "Loss", "Accuracy", "Notes"]
    data = [
        ["NormalizerAgent", "SFT 200 ex.", "0.541", "80.9%", ""],
        ["CommunicationsAgent", "SFT 80 ex.+RLVR", "0.250", "94.2%", ""],
        ["CorrelationAgent", "DSPy optimized", "—", "96.7%*", ""],
    ]
    widths = [Inches(2.8), Inches(2.4), Inches(1.2), Inches(1.4), Inches(4.6)]
    for ci, w in enumerate(widths):
        tbl.columns[ci].width = w

    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(LLM_BLUE)
        for para in cell.text_frame.paragraphs:
            para.font.bold = True
            para.font.size = Pt(13)
            para.font.color.rgb = rgb(TEXT_PRIMARY)

    for ri, row in enumerate(data, start=1):
        fill_hex = BG_CARD if ri % 2 == 1 else BG_CARD_ALT
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(fill_hex)
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(13)
                para.font.color.rgb = rgb(TEXT_PRIMARY)

    add_textbox(
        slide, Inches(0.45), Inches(3.0), Inches(12.4), Inches(0.85),
        "* Rule-based baseline scored 100% on same eval set.\n"
        "LLM advantage: confidence calibration + advisory generation.\n"
        "Advisory: LoRA when present; Ollama qwen3:8b fallback",
        font_size=11, italic=True, color_hex=TEXT_MUTED,
    )

    disc = add_rounded_rect(
        slide, Inches(0.45), Inches(4.0), Inches(12.4), Inches(2.5),
        fill_hex=BG_CARD, border_hex=GOLD, border_pt=2.0,
    )
    tf = disc.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.12)
    lines = [
        ("⚠  DISCOVERY: Reward Hacking in Communications RLVR", True, GOLD, 14),
        ("Symptom:     Model repeated boilerplate to harvest reward scores", False, TEXT_PRIMARY, 13),
        ("Root cause:  No repetition penalty in reward function", False, TEXT_PRIMARY, 13),
        ("Fix:         advisory_reward() applies unique-line ratio + length tier penalties", False, TEXT_PRIMARY, 13),
        ("Result:      Stable output; 94.2% accuracy preserved", False, TEXT_PRIMARY, 13),
    ]
    for i, (text, bold, color, size) in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = text
        para.font.size = Pt(size)
        para.font.bold = bold
        para.font.color.rgb = rgb(color)


def _info_panel(
    slide,
    left,
    top,
    width,
    height,
    header: str,
    body: str,
    *,
    border_hex: str = LLM_BLUE,
    header_hex: str = LLM_BLUE,
    body_font: str = "Calibri",
    body_size: int = 13,
) -> None:
    """Panel with colored header line."""
    box = add_rounded_rect(slide, left, top, width, height, fill_hex=BG_CARD, border_hex=border_hex)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_top = Inches(0.1)
    p0 = tf.paragraphs[0]
    p0.text = header
    p0.font.size = Pt(14)
    p0.font.bold = True
    p0.font.color.rgb = rgb(header_hex)
    p1 = tf.add_paragraph()
    p1.text = body
    p1.font.name = body_font
    p1.font.size = Pt(body_size)
    p1.font.color.rgb = rgb(TEXT_PRIMARY)


def build_slide_8_dspy(prs: Presentation) -> None:
    """Slide 8 — CorrelationAgent DSPy."""
    slide = blank_slide(prs)
    add_slide_title(slide, "CorrelationAgent — DSPy Optimized")
    add_subtitle(slide, "The decision engine behind every CONFIRMED advisory")

    lx, lw = Inches(0.5), Inches(5.8)
    rx = Inches(6.8)
    y = Inches(1.55)
    gap = Inches(0.15)

    _info_panel(
        slide, lx, y, lw, Inches(1.45),
        "What DSPy does",
        "A compiler for LLM programs — not supervised fine-tuning. DSPy optimizes "
        "prompts and few-shot examples automatically against a scoring metric, producing "
        "a compiled program that generalizes beyond its training examples.",
    )
    y += Inches(1.45) + gap
    _info_panel(
        slide, lx, y, lw, Inches(1.45),
        "Program structure",
        "Signature:  alerts + topology → root_cause + confidence\n"
        "Module:     ChainOfThought(CorrelationSignature)\n"
        "Optimizer:  BootstrapFewShot\n"
        "Metric:     confidence calibration score",
        body_font="Consolas",
        body_size=12,
    )
    y += Inches(1.45) + gap
    _info_panel(
        slide, lx, y, lw, Inches(1.2),
        "Optimized against",
        "Labeled alert clusters with known root causes. Metric rewards high confidence "
        "on true positives and penalizes overconfident wrong attributions.",
    )

    # Right — accuracy card
    add_rounded_rect(
        slide, rx, Inches(1.55), lw, Inches(1.55), fill_hex=BG_CARD, border_hex=GOLD, border_pt=2.0
    )
    add_textbox(slide, rx + Inches(0.2), Inches(1.7), lw - Inches(0.4), Inches(0.7),
                "96.7%", font_size=52, bold=True, color_hex=GOLD, align=PP_ALIGN.CENTER)
    add_textbox(slide, rx + Inches(0.2), Inches(2.35), lw - Inches(0.4), Inches(0.35),
                "correlation accuracy", font_size=14, align=PP_ALIGN.CENTER)
    add_textbox(
        slide, rx + Inches(0.15), Inches(2.7), lw - Inches(0.3), Inches(0.45),
        "Rule-based baseline: 100% on same eval set\nLLM advantage: confidence calibration",
        font_size=12, color_hex=TEXT_MUTED, align=PP_ALIGN.CENTER,
    )

    cy = Inches(3.25)
    _info_panel(
        slide, rx, cy, lw, Inches(2.05),
        "Why confidence matters",
        "Every correlate_alert() call returns a confidence score.\n"
        "That single number drives all downstream decisions:",
        border_hex=CONFIRMED_TEAL,
        header_hex=CONFIRMED_TEAL,
    )
    row_y = cy + Inches(0.95)
    add_left_accent_row(
        slide, rx + Inches(0.15), row_y, lw - Inches(0.3), Inches(0.32),
        LLM_BLUE, "< 0.85  →  PRELIMINARY advisory issued",
    )
    add_left_accent_row(
        slide, rx + Inches(0.15), row_y + Inches(0.38), lw - Inches(0.3), Inches(0.38),
        CONFIRMED_TEAL, "≥ 0.85 + count ≥ 2  →  CONFIRMED advisory issued",
    )
    add_left_accent_row(
        slide, rx + Inches(0.15), row_y + Inches(0.78), lw - Inches(0.3), Inches(0.32),
        RULE_GREEN, "health True  →  close_incident() fires",
    )

    why_box = add_rounded_rect(
        slide, rx, Inches(5.45), lw, Inches(1.05),
        fill_hex=BG_CARD_ALT, border_hex=BORDER_MUTED,
    )
    tf_w = why_box.text_frame
    tf_w.word_wrap = True
    tf_w.margin_left = Inches(0.12)
    tf_w.margin_top = Inches(0.1)
    pw0 = tf_w.paragraphs[0]
    pw0.text = "Why not fine-tune?"
    pw0.font.size = Pt(12)
    pw0.font.bold = True
    pw0.font.color.rgb = rgb(TEXT_MUTED)
    pw1 = tf_w.add_paragraph()
    pw1.text = (
        "Correlation requires reasoning over dynamic, evolving alert clusters — not "
        "pattern matching on fixed formats. DSPy's compiled few-shot approach generalizes "
        "better than SFT for this task."
    )
    pw1.font.size = Pt(12)
    pw1.font.italic = True
    pw1.font.color.rgb = rgb(TEXT_MUTED)


def _loop_step_card(
    slide,
    left,
    top,
    step: str,
    tool: str,
    kind: str,
    border_hex: str,
    kind_hex: str,
) -> None:
    """Compact step card for agentic loop slide."""
    w, h = Inches(3.5), Inches(1.0)
    card = add_rounded_rect(slide, left, top, w, h, fill_hex=BG_CARD, border_hex=border_hex)
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_top = Inches(0.06)
    p0 = tf.paragraphs[0]
    p0.text = step
    p0.font.size = Pt(12)
    p0.font.bold = True
    p0.font.color.rgb = rgb(TEXT_PRIMARY)
    p1 = tf.add_paragraph()
    p1.text = tool
    p1.font.size = Pt(11)
    p1.font.color.rgb = rgb(TEXT_MUTED)
    p2 = tf.add_paragraph()
    p2.text = kind
    p2.font.size = Pt(10)
    p2.font.italic = True
    p2.font.color.rgb = rgb(kind_hex)


def build_slide_9_agentic_loop(prs: Presentation) -> None:
    """Slide 9 — The Agentic Loop."""
    slide = blank_slide(prs)
    add_slide_title(slide, "The Agentic Loop")
    add_subtitle(slide, "ADK LlmAgent orchestrates 6 steps every 15 seconds")

    ox, oy, ow, oh = Inches(0.45), Inches(1.5), Inches(12.4), Inches(4.8)
    add_rounded_rect(
        slide, ox, oy, ow, oh, fill_hex=BG_CARD_ALT, border_hex=LLM_BLUE, border_pt=2.0
    )
    add_textbox(
        slide, ox + Inches(0.3), Inches(1.68), ow - Inches(0.6), Inches(0.4),
        "ADK LlmAgent — Orchestrator", font_size=16, bold=True, color_hex=LLM_BLUE,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide, ox + Inches(0.4), Inches(2.02), ow - Inches(0.8), Inches(0.4),
        "Conductor: reads world state, sequences tool calls, decides when to act",
        font_size=12, color_hex=TEXT_MUTED, align=PP_ALIGN.CENTER,
    )

    cw, gap = Inches(3.5), Inches(0.3)
    x0 = ox + Inches(0.55)
    row1_y, row2_y = Inches(2.5), Inches(3.65)

    cards_r1 = [
        ("① get_active_alerts()", "PrometheusMCP", "rule-based", DATA_GRAY, DATA_GRAY),
        ("② normalize_alert()", "NormalizerAgent", "LLM fine-tuned", LLM_BLUE, LLM_BLUE),
        ("③ route_alert()", "TriageAgent", "rule-based", RULE_GREEN, RULE_GREEN),
    ]
    cards_r2 = [
        ("④ correlate_alert()", "CorrelationAgent", "DSPy optimized", LLM_BLUE, LLM_BLUE),
        ("⑤ generate_advisory()", "CommunicationsAgent", "LLM fine-tuned", LLM_BLUE, LLM_BLUE),
        ("⑥ health check + close", "Prometheus + Comms", "rule-based + LLM", CONFIRMED_TEAL, CONFIRMED_TEAL),
    ]

    for i, (step, tool, kind, border, kc) in enumerate(cards_r1):
        x = x0 + i * (cw + gap)
        _loop_step_card(slide, x, row1_y, step, tool, kind, border, kc)
        if i < 2:
            add_horizontal_arrow(slide, x + cw + Inches(0.05), row1_y + Inches(0.45), Inches(0.22))

    for i, (step, tool, kind, border, kc) in enumerate(cards_r2):
        x = x0 + i * (cw + gap)
        _loop_step_card(slide, x, row2_y, step, tool, kind, border, kc)
        if i < 2:
            add_horizontal_arrow(slide, x + cw + Inches(0.05), row2_y + Inches(0.45), Inches(0.22))

    # Down arrow 3 → 4
    add_vertical_arrow(
        slide,
        x0 + 2 * (cw + gap) + cw - Inches(0.2),
        row1_y + Inches(1.05),
        Inches(0.35),
        ARROW_GRAY,
    )

    # Loop-back dashed arrow label
    add_labeled_arrow_right(
        slide,
        ox + Inches(0.8),
        row2_y + Inches(1.15),
        1.8,
        "15s cycle",
        arrow_color=LLM_BLUE,
        label_color=LLM_BLUE,
        italic=False,
        dashed=True,
    )

    callout = add_rounded_rect(
        slide, Inches(3.4), Inches(5.45), Inches(6.5), Inches(0.65),
        fill_hex=BG_CARD, border_hex=CONFIRMED_TEAL,
    )
    tf = callout.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = (
        "STEP 6 fires close_incident() only if check_service_health() returns True "
        "— the orchestrator decides, the tool executes"
    )
    p.font.size = Pt(12)
    p.font.color.rgb = rgb(TEXT_PRIMARY)
    p.alignment = PP_ALIGN.CENTER

    strip = add_rounded_rect(slide, ox, Inches(6.7), ow, Inches(0.45), fill_hex=BG_CARD)
    no_border(strip)
    tf2 = strip.text_frame
    tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
    p2 = tf2.paragraphs[0]
    p2.text = (
        "T+60–90s  first alert  →  T+2–3 cycles  CONFIRMED  →  "
        "T+1–2 cycles  SERVICE RESTORED"
    )
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(12)
    p2.font.color.rgb = rgb(TEXT_PRIMARY)


def build_slide_10_agentic_arch(prs: Presentation) -> None:
    """Slide 10 — Google ADK decisions."""
    slide = blank_slide(prs)
    add_slide_title(slide, "Google ADK LlmAgent + 8 FunctionTools")
    add_subtitle(slide, "6-step workflow  |  15s poll cycle", top=Inches(0.95), size=14)

    py = Inches(1.6)
    pw, ph = Inches(5.8), Inches(2.8)

    left = add_rounded_rect(slide, Inches(0.45), py, pw, ph, fill_hex=BG_CARD, border_hex=LLM_BLUE, border_pt=2.0)
    tf_l = left.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = Inches(0.15)
    tf_l.margin_top = Inches(0.12)
    pl0 = tf_l.paragraphs[0]
    pl0.text = "Which advisory type?"
    pl0.font.size = Pt(15)
    pl0.font.bold = True
    pl0.font.color.rgb = rgb(LLM_BLUE)
    for line in [
        "Preliminary or confirmed based on:",
        "• confidence score",
        "• alert count",
        "• flags",
    ]:
        pl = tf_l.add_paragraph()
        pl.text = line
        pl.font.size = Pt(13)
        pl.font.color.rgb = rgb(TEXT_PRIMARY)

    right = add_rounded_rect(
        slide, Inches(6.95), py, pw, ph, fill_hex=BG_CARD, border_hex=RULE_GREEN, border_pt=2.0
    )
    tf_r = right.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = Inches(0.15)
    tf_r.margin_top = Inches(0.12)
    pr0 = tf_r.paragraphs[0]
    pr0.text = "When to close?"
    pr0.font.size = Pt(15)
    pr0.font.bold = True
    pr0.font.color.rgb = rgb(RULE_GREEN)
    pr1 = tf_r.add_paragraph()
    pr1.text = (
        "Only when Prometheus health check confirms recovery — "
        "LLM coordinates, close_incident() executes"
    )
    pr1.font.size = Pt(13)
    pr1.font.color.rgb = rgb(TEXT_PRIMARY)

    div = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6.55), py + Inches(0.3), Inches(0.03), ph - Inches(0.6)
    )
    set_shape_fill(div, BORDER_MUTED)
    no_border(div)

    insight = add_rounded_rect(
        slide, Inches(0.45), Inches(4.6), Inches(12.4), Inches(1.0),
        fill_hex=BG_CARD, border_hex=GOLD, border_pt=1.5,
    )
    tf_i = insight.text_frame
    tf_i.word_wrap = True
    tf_i.vertical_anchor = MSO_ANCHOR.MIDDLE
    pi = tf_i.paragraphs[0]
    pi.text = (
        "Same alert, different outcomes — depends on evolving store state: "
        "open incidents, flags, correlation history"
    )
    pi.alignment = PP_ALIGN.CENTER
    pi.font.size = Pt(13)
    pi.font.color.rgb = rgb(TEXT_PRIMARY)

    add_textbox(slide, Inches(0.45), Inches(5.75), Inches(5.9), Inches(0.55),
                "All 6 steps orchestrated by the LLM.", font_size=13, color_hex=TEXT_MUTED)
    add_textbox(
        slide, Inches(6.5), Inches(5.75), Inches(6.3), Inches(0.7),
        "Routing and health decisions happen inside rule-based tools — not in the LLM itself.",
        font_size=13, color_hex=TEXT_MUTED,
    )


def _stat_card(slide, left, top, stat: str, label: str, caption: str, stat_color: str) -> None:
    """Results metric tile."""
    add_rounded_rect(slide, left, top, Inches(3.8), Inches(2.2), fill_hex=BG_CARD)
    add_textbox(slide, left + Inches(0.1), top + Inches(0.18), Inches(3.6), Inches(0.75),
                stat, font_size=36, bold=True, color_hex=stat_color, align=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.1), top + Inches(0.95), Inches(3.6), Inches(0.4),
                label, font_size=14, align=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.1), top + Inches(1.38), Inches(3.6), Inches(0.55),
                caption, font_size=11, color_hex=TEXT_MUTED, align=PP_ALIGN.CENTER)


def build_slide_11_results(prs: Presentation) -> None:
    """Slide 11 — Results."""
    slide = blank_slide(prs)
    add_slide_title(slide, "What It Does In Production")
    add_subtitle(slide, "Observed in demo runs — May 2026", top=Inches(0.95))

    cards = [
        ("299", "tests passing", "end-to-end, ~13s full suite", CONFIRMED_TEAL),
        ("0.90–0.95", "live confidence", "DSPyCorrelator on real alerts", LLM_BLUE),
        ("2–3 cycles", "to CONFIRMED", "after fault injection", GOLD),
        ("1–2 cycles", "to SERVICE RESTORED", "after healing", CONFIRMED_TEAL),
        ("isolated", "ad noise", "separate from valkey-cart cascade", RULE_GREEN),
        ("T+60–90s", "detection", "docker stop to first cart alert", MCP_ORANGE),
    ]
    xs = [Inches(0.45), Inches(4.55), Inches(8.65)]
    ys = [Inches(1.45), Inches(3.85)]
    for idx, (stat, label, cap, color) in enumerate(cards):
        _stat_card(slide, xs[idx % 3], ys[idx // 3], stat, label, cap, color)


def build_slide_12_live_demo(prs: Presentation) -> None:
    """Slide 12 — Live Demo."""
    slide = blank_slide(prs)
    banner = add_rounded_rect(
        slide, Inches(0.45), Inches(1.0), Inches(12.4), Inches(1.0), fill_hex=LLM_BLUE
    )
    no_border(banner)
    tf = banner.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "[ LIVE DEMO ]"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = rgb(TEXT_PRIMARY)

    add_textbox(
        slide, Inches(0.8), Inches(2.2), Inches(11.7), Inches(0.7),
        "Scenario: valkey-cart cache failure cascade\nOpenTelemetry demo environment",
        font_size=15, color_hex=TEXT_MUTED, align=PP_ALIGN.CENTER,
    )

    steps = [
        (DATA_GRAY, "Baseline", "system monitoring, ad noise isolated"),
        (ALERT_RED, "Fault", "docker stop valkey-cart"),
        (GOLD, "Detection", "cart cascade alerts, new incident opens"),
        (LLM_BLUE, "Advisory", "CONFIRMED NOC ADVISORY (confidence > 0.85)"),
        (CONFIRMED_TEAL, "Healing", "docker start valkey-cart"),
        (CONFIRMED_TEAL, "Resolution", "SERVICE RESTORED, incident removed from board"),
    ]
    card_w = Inches(9.0)
    card_x = Inches(2.15)
    y0 = Inches(3.1)
    spacing = Inches(0.68)

    for i, (accent, name, detail) in enumerate(steps):
        y = y0 + i * spacing
        add_filled_circle(slide, card_x - Inches(0.42), y + Inches(0.18), Inches(0.35), accent)
        ctf = slide.shapes[-1].text_frame
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = ctf.paragraphs[0]
        cp.text = str(i + 1)
        cp.alignment = PP_ALIGN.CENTER
        cp.font.size = Pt(12)
        cp.font.bold = True
        cp.font.color.rgb = rgb(TEXT_PRIMARY)
        add_rounded_rect(slide, card_x, y, card_w, Inches(0.62), fill_hex=BG_CARD)
        add_textbox(
            slide, card_x + Inches(0.15), y + Inches(0.1), Inches(8.5), Inches(0.5),
            f"{name}  —  {detail}", font_size=14,
        )


def build_slide_13_lessons(prs: Presentation) -> None:
    """Slide 13 — Lessons Learned."""
    slide = blank_slide(prs)
    add_slide_title(slide, "What We Learned — And Would Do Differently")

    lessons = [
        ("Baseline heuristics can beat LLM for structured tasks",
         "Establish a baseline before investing in fine-tuning"),
        ("Statelessness is a fundamental agentic challenge",
         "Design state management into architecture from day one"),
        ("Trust code over prompts for invariants",
         "Behavioral guarantees belong in code, not LLM instructions"),
        ("Noise isolation is as important as detection",
         "Understand operating environment before tuning accuracy"),
        ("Fine-tuning teaches format, not temporal awareness",
         "Know the boundary of what training can and cannot do"),
        ("Choose your orchestration framework first",
         "Build everything around it — assumptions surface only at migration"),
    ]
    c1, c2 = Inches(0.5), Inches(6.9)
    row_ys = [Inches(1.25), Inches(3.0), Inches(4.75)]
    cw, ch = Inches(5.8), Inches(1.55)

    for idx, (header, body) in enumerate(lessons):
        col = c1 if idx % 2 == 0 else c2
        top = row_ys[idx // 2]
        card = add_rounded_rect(slide, col, top, cw, ch, fill_hex=BG_CARD)
        no_border(card)
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.12)
        tf.margin_top = Inches(0.1)
        ph = tf.paragraphs[0]
        ph.text = header
        ph.font.size = Pt(13)
        ph.font.bold = True
        ph.font.color.rgb = rgb(GOLD)
        pb = tf.add_paragraph()
        pb.text = body
        pb.font.size = Pt(12)
        pb.font.color.rgb = rgb(TEXT_PRIMARY)


def build_deck() -> Presentation:
    """Assemble all 13 slides."""
    prs = new_presentation()
    build_slide_1_title(prs)
    build_slide_2_problem(prs)
    build_slide_3_system_context(prs)
    build_slide_4_architecture(prs)
    build_slide_5_tool_stack(prs)
    build_slide_6_how_it_works(prs)
    build_slide_7_models(prs)
    build_slide_8_dspy(prs)
    build_slide_9_agentic_loop(prs)
    build_slide_10_agentic_arch(prs)
    build_slide_11_results(prs)
    build_slide_12_live_demo(prs)
    build_slide_13_lessons(prs)
    return prs


def main() -> None:
    """Generate the capstone PPTX."""
    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs = build_deck()
    prs.save(str(OUTPUT_PATH))
    size_kb = OUTPUT_PATH.stat().st_size / 1024
    print(f"Created: {OUTPUT_PATH}")
    print(f"Slides: {len(prs.slides)}")
    print(f"Size: {size_kb:.1f} KB")


if __name__ == "__main__":
    main()
